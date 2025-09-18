import pandas as pd
from datetime import datetime, date
import numpy as np
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import warnings
warnings.filterwarnings("ignore")  # Ignores all warnings


def get_sheets_to_process(file_path: str, skip_first: bool = True) -> list:
    """
    Reads all sheet names from an Excel file and returns the list of sheets to process.
    
    Args:
        file_path (str): Path to the Excel file.
        skip_first (bool): If True, skips the first sheet. Default is True.
    
    Returns:
        list: List of sheet names to process.
    """
    all_sheets = pd.ExcelFile(file_path).sheet_names
    return all_sheets[1:] if skip_first else all_sheets

def read_modelwise_excel(file_path, sheet_names):
    all_data = []

    # Month mapping to standardize
    month_map = {
        "jan": "Jan", "feb": "Feb", "mar": "Mar", "apr": "Apr",
        "may": "May", "jun": "Jun", "jul": "Jul", "aug": "Aug",
        "sep": "Sep", "oct": "Oct", "nov": "Nov", "dec": "Dec",
        "march": "Mar", "april": "Apr", "june": "Jun", "july": "Jul", "august": "Aug"
    }

    for sheet in sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet)
        df.columns = df.columns.str.strip()

        # get mon and yr from sheet name 
        match = re.search(r"(jan|feb|mar|march|apr|april|may|jun|june|jul|july|aug|august|sep|oct|nov|dec)[a-z\- ]*(\d{2,4})",
                          sheet, flags=re.IGNORECASE)
        if not match:
            print(f"❌ Skipping {sheet}: cannot parse month-year")
            continue

        month_str, year_str = match.groups()
        month_std = month_map[month_str.lower()[:3]]
        year_full = year_str if len(year_str) == 4 else f"20{year_str}"
        target_mm_yy = f"{month_std}-{year_full[-2:]}"  # e.g. Nov-23

        # print(f"\nProcessing sheet {sheet} → {target_mm_yy}") #debug

        # get columns for this month + year ---
        id_cols = ["Category", "Company Name", "Model"]
        value_cols = []
        for c in df.columns:
            if c.lower().startswith(month_str.lower()) and year_full in c:
                value_cols.append(c)

        if not value_cols:
            # print(f"⚠️ Skipping {sheet}: no columns for {target_mm_yy}")
            continue
        # print(f"👉 For {sheet}, detected value_cols = {value_cols}")
        missing_ids = [c for c in id_cols if c not in df.columns]
        if missing_ids:
            # print(f"⚠️ In {sheet}, missing id cols: {missing_ids}") ## hash out continue
            continue


        # melt and pivot
        df_long = df.melt(id_vars=id_cols, value_vars=value_cols,
                          var_name="month_metric", value_name="value")
        # print(f"👉 Melted shape for {sheet}: {df_long.shape}")
        # print(df_long.head())

        df_long["metric"] = df_long["month_metric"].apply(lambda x: "Sales" if "sales" in x.lower() else "Exports")
        df_long["mm_yy"] = target_mm_yy

        df_final = df_long.pivot_table(
            index=id_cols + ["mm_yy"],
            columns="metric",
            values="value",
            aggfunc="first"
        ).reset_index()
        # print(f"👉 Pivoted shape for {sheet}: {df_final.shape}")
        # print(df_final.head())


        all_data.append(df_final)

    if not all_data:
        raise ValueError("No valid data found in the provided sheets")

    final_df = pd.concat(all_data, ignore_index=True)
    final_df.columns.name = None

    return final_df


def add_financial_quarter(df, col="mm_yy"):
    # conv mm_yy into datetime
    df["_date"] = pd.to_datetime(df[col], format="%b-%y")

    # Financial year: Apr–Mar, so add 1 if month >= Apr else year stays
    df["fy"] = df["_date"].apply(lambda d: d.year + 1 if d.month >= 4 else d.year)

    # Quarter mapping (Apr–Jun → Q1, Jul–Sep → Q2, Oct–Dec → Q3, Jan–Mar → Q4)
    def get_quarter(d):
        if d.month in [4,5,6]:
            return "Q1"
        elif d.month in [7,8,9]:
            return "Q2"
        elif d.month in [10,11,12]:
            return "Q3"
        else:
            return "Q4"

    df["quarter"] = df["_date"].apply(get_quarter)

    # Build final qtr_yy col
    df["qtr_yy"] = df["quarter"] + "-" + df["fy"].astype(str).str[-2:]

    # Drop helper cols
    df = df.drop(columns=["_date","fy","quarter"])
    return df


company_df_dict = {
    'Piaggio Vehicles Pvt Ltd': 'df_piaggio',
    'TVS Motor Company Ltd': 'df_tvs',
    'Hero MotoCorp Ltd': 'df_hero',
    'Honda Motorcycle & Scooter India Pvt Ltd': 'df_honda',
    'India Yamaha Motor Pvt Ltd': 'df_yamaha',
    'Suzuki Motorcycle India Pvt Ltd': 'df_suzuki',
    'Chetak Technology Ltd': 'df_chetak',
    'Okinawa Autotech Pvt. Ltd': 'df_okinawa',
    'Ather Energy Pvt. Ltd': 'df_ather',
    'Bajaj Auto Ltd': 'df_bajaj',
    'India Kawasaki Motors Pvt Ltd': 'df_kawasaki',
    'Triumph Motorcycles India Pvt Ltd': 'df_triumph',
    'Mahindra Two Wheelers Ltd': 'df_mahindra',
    'Royal-Enfield (Unit of Eicher Motors)': 'df_royal_enfield'
}


model_filters = {
    "df_bajaj":[
       'Avenger', 'Boxer', 'CT', 'CT 150', 'Discover',
       'Dominar', 'Freedom', 'Husqvarna', 'KTM', 'Platina', 'Pulsar',
       'Triumph'
    ],
    "df_tvs": [
        "Star", "Apache", "Sport", "Raider", "BMW", "Ronin", "RTR 310", "RR 310", "Radeon"
    ],
    "df_royal_enfield": [
        "Classic 350", "Meteor 350", "Hunter 350", "Himalayan", "Bullet 350",
        "Guerrilla", "Single Cylinder Total", "Super Meteor", "650 Twins",
        "Shotgun", "Twin Cylinder Total"
    ],
    "df_yamaha": [
        "FZ / Fazer", "Saluto RX", "Gladiator/ Saluto", "MT", "SZ",
        "FZ 25", "YZF R15", "CRUX"
    ],
    "df_suzuki": [
        "Gixxer 150 Series", "Gixxer 250", "V-Strom SX", "Hayate"
    ],
    "df_honda": [
        "Shine", "Dream Series", "Hornet 2.0", "X Blade", "Unicorn",
        "Livo", "SP 160", "CB 350", "CB 160R", "CB 200X", "H'ness",
        "CB Twister", "MC 300N", "CB 300F"
    ],
    "df_hero": [
        "HUNK", "HF DELUXE", "X PULSE", "SPLENDOR", "Xtreme",
        "GLAMOUR", "PASSION", "HF DAWN", "ACHIEVER", "Karizma"
    ],
    "df_piaggio": [
        "RS", "457 CC", "Tuono 457 cc"
    ]
}


brand_short_names = {
    "df_bajaj": "Bajaj+KTM+Triumph",
    "df_tvs": "TVS",
    "df_royal_enfield": "Royal Enfield",
    "df_yamaha": "Yamaha",
    "df_suzuki": "Suzuki",
    "df_honda": "Honda",
    "df_hero": "Hero",
    "df_piaggio": "Piaggio"
}


def filter_quarters(df):
    # Keep only columns with Qx-YY format or 'Model'
    quarter_cols = [col for col in df.columns if re.match(r'Q\d-\d{2}', col)]
    
    # Extract FY year from column
    fy_years = [int(col.split('-')[1]) for col in quarter_cols]
    
    latest_fy = max(fy_years)
    prev_fy = latest_fy - 1
    
    # Filter columns for latest FY and previous FY
    selected_cols = [col for col in quarter_cols if int(col.split('-')[1]) in [prev_fy, latest_fy]]
    
    # Sort columns: first by FY, then by quarter number
    def col_sort_key(col):
        q, yy = col.split('-')
        return (int(yy), int(q[1]))
    
    selected_cols_sorted = sorted(selected_cols, key=col_sort_key)
    
    # Include Model column at front
    final_cols = ['Model'] + selected_cols_sorted
    
    return df[final_cols].copy()


def add_grand_total_row(df: pd.DataFrame, model_col: str = "Model") -> pd.DataFrame:
    """
    Add a 'Grand Total' row at the bottom:
    - model_col will have 'Grand Total'
    - numeric columns will be summed vertically
    - other non-numeric columns (except model_col) will be NaN
    """
    out = df.copy()
    
    # Prepare dict for new row
    total_row = {}
    for col in out.columns:
        if col == model_col:
            total_row[col] = "Grand Total"
        elif pd.api.types.is_numeric_dtype(out[col]):
            total_row[col] = out[col].sum(skipna=True)
        else:
            total_row[col] = pd.NA  # keep as NaN for non-numeric

    # Append new row
    out = pd.concat([out, pd.DataFrame([total_row])], ignore_index=True)
    
    return out


def add_fy_mean_next_to_q4(df: pd.DataFrame) -> pd.DataFrame:
    """
    Calculate rounded mean of Q1–Q4 of the lower FY, treating NaNs as 0 for calculation,
    and insert as 'FY-yy' immediately after Q4 of that FY.
    Original NaNs in df are preserved.
    """
    out = df.copy()
    cols = out.columns.tolist()

    # Detect fiscal years from Q1 and Q4 columns
    fy_quarters = [c for c in cols if c.startswith("Q1-") or c.startswith("Q4-")]
    fy_years = sorted({c.split("-")[1] for c in fy_quarters}, key=int)

    if not fy_years:
        raise ValueError("No fiscal years found in DataFrame")

    # Pick the lower FY
    lower_fy = fy_years[0]

    # Quarters of lower FY
    q_cols = [f"Q{i}-{lower_fy}" for i in range(1, 5)]

    # Ensure all quarters exist
    if not all(q in out.columns for q in q_cols):
        raise ValueError(f"Missing one or more quarters for FY-{lower_fy}")

    # Temporarily fill NaNs with 0 for mean calculation
    fy_mean = out[q_cols].fillna(0).mean(axis=1).round().astype("Int64")

    # Insert FY mean after Q4-lower
    q4_idx = out.columns.get_loc(f"Q4-{lower_fy}")
    out.insert(q4_idx + 1, f"FY-{lower_fy}", fy_mean)

    return out


def add_growth_metrics(df: pd.DataFrame) -> pd.DataFrame:
    """
    Add YoY and QoQ growth metrics to filtered FY quarters DataFrame.

    Expects df to have quarter columns like: 'Q1-25', 'Q2-25', 'Q3-25', 'Q4-25', 'Q1-26'
    """
    out = df.copy()
    cols = df.columns.tolist()

    # Detect available fiscal years (just last part after "-")
    fy_quarters = [c for c in cols if c.startswith("Q1-") or c.startswith("Q4-")]
    fy_years = sorted({c.split("-")[1] for c in fy_quarters}, key=int)

    if len(fy_years) < 2:
        raise ValueError("Need at least two FYs (current + next) to compute growth")

    lower_fy, upper_fy = fy_years[0], fy_years[1]

    # Column names we need
    lower_q1 = f"Q1-{lower_fy}"
    lower_q4 = f"Q4-{lower_fy}"
    upper_q1 = f"Q1-{upper_fy}"

    # calculate YoY Growth %
    if lower_q1 in out.columns and upper_q1 in out.columns:
        out["YoY Gr%"] = round(
            100 * (out[upper_q1] - out[lower_q1]) / out[lower_q1], 0
        ).astype("Int64")

    # calculate QoQ Growth %
    if lower_q4 in out.columns and upper_q1 in out.columns:
        out["QoQ Gr%"] = round(
            100 * (out[upper_q1] - out[lower_q4]) / out[lower_q4], 0
        ).astype("Int64")

    return out


def replace_large_numbers(df, threshold=1e8):
    """
    Replace numeric values in the dataframe exceeding ±threshold with NaN.
    
    Args:
        df (pd.DataFrame): Input dataframe
        threshold (float): Threshold beyond which values are treated as invalid
    
    Returns:
        pd.DataFrame: DataFrame with large numeric values replaced by NaN
    """
    df_copy = df.copy()
    num_cols = df_copy.select_dtypes(include=["number"]).columns
    df_copy[num_cols] = df_copy[num_cols].applymap(
        lambda x: np.nan if abs(x) > threshold else x
    )
    return df_copy


def filter_company_models(company_dfs, model_filters, brand_short_names):
    """
    Filter company DataFrames using model_filters.
    Only keep companies which have a model filter defined.
    """
    filtered_brands_models = {}

    for df_name, models in model_filters.items():
        if df_name not in company_dfs:
            continue  # skip companies not in pivoted dfs

        df = company_dfs[df_name].copy()
        allowed = [m.strip().lower() for m in models if m]

        # Filter by allowed models
        filtered_df = df[df["Model"].str.strip().str.lower().isin(allowed)].copy()
        if filtered_df.empty:
            continue  # skip if nothing left after filtering

        filtered_brands_models[df_name] = {
            "brand": brand_short_names.get(df_name, df_name),
            "data": filtered_df
        }

    return filtered_brands_models



def preprocess_royal_enfield(df_royal_enfield: pd.DataFrame) -> pd.DataFrame:
    """
    Preprocess Royal Enfield DataFrame:
    - Separate into single-cylinder and twin-cylinder models.
    - Add 'Single Cylinder Total' and 'Twin Cylinder Total'.
    - Add 'Grand Total' at the end (ignoring the two subtotal rows).
    """

    # Define groups
    single_models = [
        "Classic 350", "Meteor 350", "Hunter 350", "Himalayan", "Bullet 350", "Guerrilla"
    ]
    twin_models = [
        "Super Meteor", "650 Twins", "Shotgun"
    ]

    df = df_royal_enfield.copy()
    df["Model"] = df["Model"].str.strip()

    # --- Single-cylinder block ---
    single_df = df[df["Model"].str.lower().isin([m.lower() for m in single_models])].copy()
    if not single_df.empty:
        single_sum = single_df.select_dtypes(include="number").sum()
        single_total = pd.DataFrame([{**{"Model": "Single Cylinder Total"}, **single_sum.to_dict()}])
        single_df = pd.concat([single_df, single_total], ignore_index=True)

    # --- Twin-cylinder block ---
    twin_df = df[df["Model"].str.lower().isin([m.lower() for m in twin_models])].copy()
    if not twin_df.empty:
        twin_sum = twin_df.select_dtypes(include="number").sum()
        twin_total = pd.DataFrame([{**{"Model": "Twin Cylinder Total"}, **twin_sum.to_dict()}])
        twin_df = pd.concat([twin_df, twin_total], ignore_index=True)

    # --- Combine blocks ---
    final_df = pd.concat([single_df, twin_df], ignore_index=True)

    # --- Grand Total (ignore subtotal rows) ---
    base_df = final_df[~final_df["Model"].isin(["Single Cylinder Total", "Twin Cylinder Total"])]
    if not base_df.empty:
        grand_sum = base_df.select_dtypes(include="number").sum()
        grand_row = pd.DataFrame([{**{"Model": "Grand Total"}, **grand_sum.to_dict()}])
        final_df = pd.concat([final_df, grand_row], ignore_index=True)

    return final_df


def dfs_to_ppt(dfs: list, titles: list, filename: str, footer: str = "Source : SIAM"):
    """
    Export multiple pandas DataFrames to a PowerPoint file.
    Each DataFrame gets its own slide with:
      - Header image with title inside it
      - Formatted table
      - Footer
    """

    prs = Presentation()

    #slide size: 34 cm wide × 19 cm high
    prs.slide_width = Cm(34)
    prs.slide_height = Cm(19)

    slide_layout = prs.slide_layouts[6]  # Blank layout

    for df, title in zip(dfs, titles):
        slide = prs.slides.add_slide(slide_layout)

        #  Header image across full slide
        header_img_path = "img/header.png"
        header_left, header_top = Cm(0), Cm(0)
        header_width = prs.slide_width  # full slide width
        slide.shapes.add_picture(header_img_path, header_left, header_top, width=header_width)

        #  Title inside header 
        txBox = slide.shapes.add_textbox(Cm(0.5), Cm(0.2), Cm(20), Cm(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)  # White
        run.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT

        # Table slide 
        rows, cols = df.shape[0] + 1, df.shape[1]
        margin = Cm(0.5)           # left/right margin
        header_height = Cm(2.5)    # approx height of header image
        gap = Cm(1)                # gap between header and table

        left = margin
        top = header_height + gap
        width = prs.slide_width - 2 * margin
        height = prs.slide_height - (header_height + gap + Cm(2))  # leave space for footer
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Equal column widths
        col_width = width / cols
        for j in range(cols):
            table.columns[j].width = int(col_width)

        # Header Row
        for j, col_name in enumerate(df.columns):
            cell = table.cell(0, j)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
            run.font.name = "Calibri"
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # --- Data Rows ---
        for i in range(df.shape[0]):
            for j, col_name in enumerate(df.columns):
                val = df.iat[i, j]
                cell = table.cell(i + 1, j)

                # Format values
                if col_name in ["YoY Gr%", "QoQ Gr%"]:
                    if pd.isna(val):
                        cell.text = "--"
                    else:
                        cell.text = f"{val}%"
                        run = cell.text_frame.paragraphs[0].runs[0]
                        if val < 0:
                            run.font.color.rgb = RGBColor(255, 0, 0)  # Red
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                elif isinstance(val, (int, float)):
                    if pd.isna(val):
                        cell.text = "--"
                    else:
                        cell.text = f"{int(val):,}"
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else:
                    cell.text = str(val)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                # Font formatting
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(14)
                run.font.name = "Calibri"

                # Make first column (Model) bold
                if j == 0:
                    run.font.bold = True

        #  Last row formatting (Bold, dark blue) ---
        last_row_idx = df.shape[0]
        for j in range(cols):
            cell = table.cell(last_row_idx, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.name = "Calibri"

        #
        footer_left, footer_top = margin, prs.slide_height - Cm(1.5)
        txBox = slide.shapes.add_textbox(footer_left, footer_top, width, Cm(1))
        tf = txBox.text_frame
        tf.text = footer
        tf.paragraphs[0].runs[0].font.size = Pt(9)
        tf.paragraphs[0].runs[0].font.name = "Calibri"

    prs.save(filename)
    # print(f"Saved PPT with {len(dfs)} slides: {filename}")






def summarize_by_quarter(df: pd.DataFrame) -> pd.DataFrame:
    return (
        df.groupby(["Company Name", "Model", "qtr_yy"], as_index=False)[["Sales", "Exports"]]
        .mean()
    )

def create_pivot(df: pd.DataFrame, values: str = "Exports") -> pd.DataFrame:
    df_pivot = (
        df.pivot_table(
            index=["Company Name", "Model"],
            columns="qtr_yy",
            values=values,
            aggfunc="mean"
        )
        .reset_index()
    )
    num_cols = df_pivot.select_dtypes(include="number").columns
    df_pivot[num_cols] = df_pivot[num_cols].round(0).astype("Int64")
    return df_pivot


def split_companies(df_pivot: pd.DataFrame, company_df_dict: dict, filter_quarters_func) -> dict:
    result = {}
    for company, df_name in company_df_dict.items():
        temp = df_pivot[df_pivot["Company Name"] == company].copy()
        temp = filter_quarters_func(temp)   # using your defined function
        result[df_name] = temp
    return result

def apply_grand_total(filtered_dict: dict, add_grand_total_func) -> dict:
    for key, value in filtered_dict.items():
        value["data"] = add_grand_total_func(value["data"])
    return filtered_dict



# Step 12 → Special Royal Enfield preprocessing
def apply_royal_enfield_patch(filtered_dict: dict, preprocess_func) -> dict:
    if "df_royal_enfield" in filtered_dict:
        filtered_dict["df_royal_enfield"]["data"] = preprocess_func(
            filtered_dict["df_royal_enfield"]["data"]
        )
    return filtered_dict

# Step 13–15 → Apply multiple transforms in sequence
def apply_transformations(filtered_dict: dict,
                          fy_mean_func,
                          growth_func,
                          replace_func) -> dict:
    for key, value in filtered_dict.items():
        df = value["data"]
        df = fy_mean_func(df)
        df = growth_func(df)
        df = replace_func(df)
        value["data"] = df
    return filtered_dict

def process_workflow(file_path, sheets_to_process,
                     company_df_dict, model_filters, brand_short_names,
                     filter_quarters,
                     filter_company_models,
                     add_grand_total_row,
                     preprocess_royal_enfield,
                     add_fy_mean_next_to_q4,
                     add_growth_metrics,
                     replace_large_numbers,
                     dfs_to_ppt,
                     values="Exports",
                     output_ppt="brands_exports_final_orch_og.pptx"):

    # Step 1–2
    df_final = read_modelwise_excel(file_path, sheets_to_process)
    df_final = add_financial_quarter(df_final)
    # print("df_final:")
    # print(df_final.head())
    # print("=="*30)

    # Step 3–4
    df_qtr_summary = summarize_by_quarter(df_final)
    # print("df_qtr_summary:")
    # print(df_qtr_summary.head())
    df_pivot = create_pivot(df_qtr_summary, values=values)
    # print("df_pivot:")
    # print(df_pivot.head())
    # print("=="*30)

    # Step 6–7
    company_dfs = split_companies(df_pivot, company_df_dict, filter_quarters)
    # print("company_dfs:")
    # print(company_dfs)
    # print("=="*30)

    # Step 10
    filtered_brands_models = filter_company_models(company_dfs, model_filters, brand_short_names)
    # print("filtered_brands_models:")
    # print(filtered_brands_models)
    # print("=="*30)

    # Step 11–12–13–14–15
    filtered_brands_models = apply_grand_total(filtered_brands_models, add_grand_total_row)
    filtered_brands_models = apply_royal_enfield_patch(filtered_brands_models, preprocess_royal_enfield)
    filtered_brands_models = apply_transformations(filtered_brands_models,
                                                  add_fy_mean_next_to_q4,
                                                  add_growth_metrics,
                                                  replace_large_numbers)
    # print("filtered_brands_models_updated:")
    # print(filtered_brands_models)
    # print("=="*30)

    # Step 16 → Collect dfs + titles
    dfs = [v["data"] for v in filtered_brands_models.values()]
    # print("dfs:")
    # print(dfs)
    # print("=="*30)
    titles = [v["brand"] for v in filtered_brands_models.values()]
    # print("titles:")
    # print(titles)
    # print("=="*30)
    

    # Step 17 → Save PPT
    dfs_to_ppt(dfs, titles, filename=output_ppt)
    # print("Completed")

    return output_ppt
